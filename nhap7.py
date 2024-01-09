from pptx import *
import win32com
from win32com import client
# import aspose.slides as slidess
try:
    from StringIO import StringIO ## for Python 2
except ImportError:
    from io import StringIO
def delete_allslide(filename):
    #for j,da in enumerate(xml_slides):
        #print(j)
    #filename = '10slide.pptx'
    prs = Presentation(filename)

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    print(slides)
    for i in slides:

        xml_slides.remove(i)
    # for i in slides:
    #     xml_slides.remove(i)
    prs.save('1slide.pptx')

def moveshapeto_allslide(filename):
    #for j,da in enumerate(xml_slides):
        #print(j)
    #filename = '10slide.pptx'
    prs = Presentation(filename)

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    print(slides)
    for slide in prs.slides:
        for shape in slide.shapes:
            shape.left = -10000
            shape.top = -10000
    # for i in slides:
    #     xml_slides.remove(i)
    prs.save('1slide.pptx')
SLD_LAYOUT_TITLE_AND_CONTENT = 1
prs = Presentation(r'C:\Users\Admin\anaconda3\envs\backuptamthoi3\pptx2h5p_1_2_deep_copy_addimages\10slide.pptx'
                    )
# b= pptx.slide.Slide
# for i in a.slides
#
# prs = Presentation(filename)

xml_slides = prs.slides._sldIdLst
for i,j in enumerate(prs.slides):
     print(j)
# for shape in slide.placeholders:
#      print('%d %s' % (shape.placeholder_format.idx, shape.name))
#delete text
# textbox = shapes[textbox_idx]
# sp = textbox.element
# sp.getparent().remove(sp)
# for slide in prs.slides:
#    for shape in slide.shapes:
#        shape.left = -900000000
#        shape.top = -90000000
#from pptx import Presentation
pres = Presentation('10slide.pptx')
# for slide in pres.slides:
#     for shape in slide.shapes:
#         image = shape.image
#         blob = image.blob
#         ext = image.ext
#         with open(f'image.{ext}', 'wb') as file:
#             file.write('tam\%s'%('')+blob)

slide = pres.slides[0]
shape = slide.shapes[0]
image = shape.image
blob = image.blob
ext = image.ext
with open(f'image.{ext}', 'wb') as file:
    file.write('tam\%s' % ('') + blob)
# prs.save('1slide.pptx')