from pptx import Presentation
import os
import shutil
from pptx2h5p import ppt2png
#prs = Presentation()
# prs.save('test.pptx')
# os.remove('test.pptx')
#from pptx import Presentation
import copy
import six
# name of your pptx
def getbackgroundofslide(linkfile):
    newlinkfile=linkfile[:-5]+str('_cp.pptx')
    print(newlinkfile)
    #new_title=newlinkfile
    shutil.copy(linkfile, newlinkfile)
    prs = Presentation(newlinkfile)
    # prs = Presentation(path_to_presentation)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    for slide in prs.slides:

        for shapea in slide.shapes:
            # name of your pptx
            # or whatever slide your pictures on
            pic = shapea  # select the shape holding your picture and you want to remove

            pic = pic._element
            pic.getparent().remove(pic)  # delete the shape with the picture

                    #     # for j in i:
    #     #     print(j)
    prs.save(newlinkfile)  # save changes
    return newlinkfile#,new_title

#getbackgroundofslide(r'C:\Users\Admin\anaconda3\envs\backuptamthoi3\pptx2h5p_1_2_deep_copy_addimages\10slide.pptx')
linkfile=r'10slide.pptx'
newlinkfile=linkfile[:-5]+str('_cp.pptx')
#print(newlinkfile)
# newlinkfile=getbackgroundofslide(linkfile)
# # #new_title=newlinkfile
#shutil.copy(linkfile, newlinkfile)
# prs = Presentation(newlinkfile)
#getbackgroundofslide(linkfile)
#os.remove(newlinkfile)
