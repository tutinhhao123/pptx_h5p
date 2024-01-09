import json
import os
import sys
import uuid
import io
from pptx import Presentation
from copy import deepcopy
import nhap4
import shutil
from nhap4 import traveobj_slide,convert_donvi_text_xycxcy_2_h5p
import json
from zipfile import ZipFile
from natsort import natsorted
import edit_colon2_
from win32com import client
from get_image_size import get_image_size
path_file_ptpx_rar=r""
VERSION = "1.2"
import shutil
YEAR = "2021"
target_ratio = 2  # target aspect ratio for slides in h5p
#from imread_XML_infolder import load_xml_from_folder
from imread_XML_infolder import *

from edit_string_pptx2rar import edit_name_pptx2rar
#import xml2json.json_object
#import xml2json.xml2json as xml2json
#import xml2json.elementree_withedittext
import copy_elementree_withedittext
import elementtre_color as elementree_withedittextt
class Create_obj_p():          # leave this empty
    def __init__(self):   # constructor function using self
        self.text = None  # variable using self.
        self.color = None  # variable using self
def delete_text_element1_inslide(content):
    for i, a in enumerate(content["presentation"]["slides"]):
        print('slide', i)
        # print( a)
        print(count_keys(a['elements']))
        for j, elements in enumerate(a['elements']):
            # if(j>0):
            print(j)
            if j > 0 and j < 2:
                del a['elements'][j]
    return content


def getbackgroundofslide(linkfile):
    newlinkfile=linkfile[:-5]+str('_cp.pptx')
    print(newlinkfile)
    new_title=newlinkfile
    shutil.copy(linkfile, newlinkfile)
    prs = Presentation(newlinkfile)
    # prs = Presentation(path_to_presentation)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    # for slide in prs.slides:
    #     for shape in slide.shapes:
    #         # name of your pptx
    #         # or whatever slide your pictures on
    #         pic = shape  # select the shape holding your picture and you want to remove
    #
    #         pic = pic._element
    #         pic.getparent().remove(pic)  # delete the shape with the picture

    #     # for j in i:
    #     #     print(j)
    for slide in prs.slides:
        for shape in slide.shapes:
            shape.left = -9000000
            shape.top = -9000000
    prs.save(newlinkfile)  # save changes
    return newlinkfile#,new_title

def dithang_shape_object_ver001(arrayobj,arr_none):
    #print('#-----di thang obj ver001-----------')
    #r"""\%s""" % ("")
    #x,y,width,height,text,subContentId
    #x,y,width,height= number -> nothing
    #text,subcontent is string -> ""
    print('#-------------di thang obj-------------------')
    listobj_textfield=[]
    listobj_text_slide=[]

    list_none=arr_none
    for i,x in enumerate(arrayobj):
        listobj_shapefield = []
        if not x:
            print('shape ',i,' in slide None')
            #listobj_text_slide.append('')
        else:
            print('slide truoc add', i)
            for j,z in enumerate(x):
            #cpmverttext_h5p
            #print(j, dataa.img, dataa.rId, dataa.prst)
            #print('tesst dataa',type(dataa.x),dataa.y,dataa.cx,dataa.cy)
                print('dithang_shape_object_ver001', z.x, z.y, z.cx, z.cy,type(z.x))
                print('check color',z.img,z.prst,z.color)
           # print('check type xshape',type(z.x))
           # z.x, z.y, z.cx, z.cy = convert_donvi_text_xycxcy_2_h5p(z.x, z.y, z.cx, z.cy)
#                z.x, z.y, z.cx, z.cy = convert_donvi_text_xycxcy_2_h5p(z.x, z.y, z.cx, z.cy)
#                 if(z.x[0] == '-'):
#                     z.x = str(0)
#                 if(z.y[0] == '-'):
#                     z.y = str(0)
#                 if(z.cx[0] == '-'):
#                     z.cx = str(0)
#                 if(z.cy[0] == '-'):
#                     z.cy = str(0)

            # if  z.x is None or  or z.cx is None or z.cy is None:
            #     print('nulllllllllllllllllllllllllllllllllll')
                if z.y is None:
                    print('nulllllllllllllllllllllllllllllllllll')
                if z.x is None:
                    print('nulllllllllllllllllllllllllllllllllll')
                if z.cy is None:
                    print('nulllllllllllllllllllllllllllllllllll')
                if z.cx is None:
                    print('nulllllllllllllllllllllllllllllllllll')
                #xxx = int(z.x)
                #print('checknull',z.y,z.cx,z.cy,z.color,z.prst,z.img)
                    #
                # z.x = str((100.0 * int(z.x)) / 9144000.0)
                #
                # z.y = str((100.0 * int(z.y)) / 6858000.0)
                # z.cx = str((100.0 * int(z.cx)) / 9144000.0)
                # z.cy = str((100.0 * int(z.cy)) / 6858000.0)
            #if True:
                #h5p images
                if True:


                    if z.img !='':
                        if z.img[-3:]=='gif':
                            h5p_gif = {
                          "x": z.x,
                          "y": z.y,
                          "width": z.cx,
                          "height": z.cy,
                          "action": {
                            "library": "H5P.Image 1.1",
                            "params": {
                              "decorative": False,
                              "contentName": "Image",
                              "file": {
                                "path": z.img,
                                "mime": "image/gif",
                                "copyright": {
                                  "license": "U"
                                },
                                "width": z.cx,
                                "height": z.cy
                              }
                            },
                            "subContentId": str(nhap4.uuid.uuid4()),
                            "metadata": {
                              "contentType": "Image",
                              "license": "U",
                              "title": "Untitled Image",
                              "authors": [],
                              "changes": []
                            }
                          },
                          "alwaysDisplayComments": False,
                          "backgroundOpacity": 0,
                          "displayAsButton": False,
                          "buttonSize": "big",
                          "goToSlideType": "specified",
                          "invisible": False,
                          "solution": ""
                        }
                            listobj_shapefield.append(h5p_gif)
                        if z.img[-3:] != 'gif': #and z.img[-3:] =="png":
                            if z.img[-3:]=='jpeg' or z.img[-3:]=='png' or z.img[-3:]=='PNG' or z.img[-3:]=='jpg':
                        # lấy ảnh khác png nữa để lấy full ảnh.
                        #if z.img[-3:] != 'gif' and z.img[-3:] == "png":
                                h5p_img_mini = {
                                    "x": z.x,
                                    "y": z.y,
                                    "width": z.cx,
                                    "height": z.cy,
                                    "action": {
                                        "library": "H5P.Image 1.1",
                                        "params": {
                                            "decorative": False,
                                            "contentName": "Image",
                                            "file": {
                                                "path": z.img,
                                                "mime": "image/gif",
                                                "copyright": {
                                                    "license": "U"
                                                },
                                                "width": z.cx,
                                                "height": z.cy
                                            }
                                        },
                                        "subContentId":  str(nhap4.uuid.uuid4()),
                                        "metadata": {
                                            "contentType": "Image",
                                            "license": "U",
                                            "title": "Untitled Image",
                                            "authors": [],
                                            "changes": []
                                        }
                                    },
                                    "alwaysDisplayComments": False,
                                    "backgroundOpacity": 0,
                                    "displayAsButton": False,
                                    "buttonSize": "big",
                                    "goToSlideType": "specified",
                                    "invisible": False,
                                    "solution": ""
                                }
                                listobj_shapefield.append(h5p_img_mini)
                    #h5p shape rectanggle
                    if z.prst!='':
                        if z.prst=='rect':
                            if  z.color=='':

                                z.color='ffb3ff'
                            h5ptext_rectangle={
              "x": z.x,
              "y": z.y,
              "width": z.cx,
              "height": z.cy,
              "action": {
                "library": "H5P.Shape 1.0",
                "params": {
                  "shape": {
                    "fillColor": "#%s" % (str(z.color)),
                    "borderWidth": 0,
                    "borderStyle": "solid",
                    "borderColor": "#000",
                    "borderRadius": 0
                  },
                  "type": "rectangle",
                  "line": {
                    "borderColor": "#000000",
                    "borderWidth": "1",
                    "borderStyle": "solid"
                  }
                },
                "subContentId":  str(nhap4.uuid.uuid4()),
                "metadata": {
                  "contentType": "Shapes"
                }
              },
              "backgroundOpacity": 0,
              "displayAsButton": False,
              "buttonSize": "big",
              "goToSlideType": "specified",
              "invisible": False
            }
                            listobj_shapefield.append(h5ptext_rectangle)
                        if z.prst=='roundRect' :
                            if  z.color=='':

                                z.color='ffb3ff'
                            h5ptext_roundRect = {
                                "x": z.x,
                                "y": z.y,
                                "width": z.cx,
                                "height": z.cy,
                                "action": {
                                    "library": "H5P.Shape 1.0",
                                    "params": {
                                        "shape": {
                                            "fillColor": "#%s" % (str(z.color)),
                                            "borderWidth": 0,
                                            "borderStyle": "solid",
                                            "borderColor": "#000",
                                            "borderRadius": 3
                                        },
                                        "type": "rectangle",
                                        "line": {
                                            "borderColor": "#000000",
                                            "borderWidth": "1",
                                            "borderStyle": "solid"
                                        }
                                    },
                                    "subContentId": str(nhap4.uuid.uuid4()),
                                    "metadata": {
                                        "contentType": "Shapes"
                                    }
                                },
                                "backgroundOpacity": 0,
                                "displayAsButton": False,
                                "buttonSize": "big",
                                "goToSlideType": "specified",
                                "invisible": False
                            }
                            listobj_shapefield.append(h5ptext_roundRect)
                        if z.prst == 'line':
                            if  z.color=='':

                                z.color='ffb3ff'
                            if z.cx==0:
                                # h5p shape vertical-line - theo chieu doc
                                h5p_vertical_line = {
                                "x": z.x,
                                "y": z.y,
                                "width": 10,
                                "height": z.cy,
                                    "action": {
                                        "library": "H5P.Shape 1.0",
                                        "params": {
                                            "line": {
                                                "borderWidth": 1,
                                                "borderStyle": "solid",
                                                "borderColor": "#000"
                                            },
                                            "type": "vertical-line",
                                            "shape": {
                                                "fillColor": "#%s" % (str(z.color)),
                                                "borderColor": "#000000",
                                                "borderWidth": "0",
                                                "borderStyle": "solid",
                                                "borderRadius": "0"
                                            }
                                        },
                                        "subContentId": str(nhap4.uuid.uuid4()),
                                        "metadata": {
                                            "contentType": "Shapes"
                                        }
                                    },
                                    "backgroundOpacity": 0,
                                    "displayAsButton": False,
                                    "buttonSize": "big",
                                    "goToSlideType": "specified",
                                    "invisible": False
                                }
                                listobj_shapefield.append(h5p_vertical_line)
                            if z.cy == 0:
                                 # h5p shape horizontal-line - theo chieu ngang
                                h5p_horizontal_line = {
                                    "x": z.x,
                                "y": z.y,
                                "width": z.cx,
                                "height": 10,
                                    "action": {
                                        "library": "H5P.Shape 1.0",
                                        "params": {
                                            "line": {
                                                "borderWidth": 1,
                                                "borderStyle": "solid",
                                                "borderColor": "#000"
                                            },
                                            "type": "horizontal-line",
                                            "shape": {
                                                "fillColor": "#%s" % (str(z.color)),
                                                "borderColor": "#000000",
                                                "borderWidth": "0",
                                                "borderStyle": "solid",
                                                "borderRadius": "0"
                                            }
                                        },
                                        "subContentId": str(nhap4.uuid.uuid4()),
                                        "metadata": {
                                            "contentType": "Shapes"
                                        }
                                    },
                                    "backgroundOpacity": 0,
                                    "displayAsButton": False,
                                    "buttonSize": "big",
                                    "goToSlideType": "specified",
                                    "invisible": False
                                }
                                listobj_shapefield.append(h5p_horizontal_line)

                        if z.prst == 'circle':
                            if  z.color=='':

                                z.color='ffb3ff'
                            h5p_cirle = {
                                    "x": z.x,
                                "y": z.y,
                                "width": z.cx,
                                "height": z.cy,
                                "action": {
                                    "library": "H5P.Shape 1.0",
                                    "params": {
                                        "shape": {
                                            "fillColor": "#%s"%(str(z.color)),
                                            "borderWidth": 0,
                                            "borderStyle": "solid",
                                            "borderColor": "#000",
                                            "borderRadius": "0"
                                        },
                                        "type": "circle",
                                        "line": {
                                            "borderColor": "#000000",
                                            "borderWidth": "1",
                                            "borderStyle": "solid"
                                        }
                                    },
                                    "subContentId":  str(nhap4.uuid.uuid4()),
                                    "metadata": {
                                        "contentType": "Shapes"
                                    }
                                },
                                "backgroundOpacity": 0,
                                "displayAsButton": False,
                                "buttonSize": "big",
                                "goToSlideType": "specified",
                                "invisible": False
                            }
                            listobj_shapefield.append(h5p_cirle)



                    #add theo slide
                    #listobj_text_slide.append(listobj_shapefield)

            if False:
                    #h5p shape cirle
                    h5p_cirle={
                  "x": 23.04147465437788,
                  "y": 32.511126933192884,
                  "width": 14.09,
                  "height": 27.831977000000002,
                  "action": {
                    "library": "H5P.Shape 1.0",
                    "params": {
                      "shape": {
                        "fillColor": "#fff",
                        "borderWidth": 0,
                        "borderStyle": "solid",
                        "borderColor": "#000",
                        "borderRadius": "0"
                      },
                      "type": "circle",
                      "line": {
                        "borderColor": "#000000",
                        "borderWidth": "1",
                        "borderStyle": "solid"
                      }
                    },
                    "subContentId": "5685a3b6-3369-4089-853c-da42eb7fa791",
                    "metadata": {
                      "contentType": "Shapes"
                    }
                  },
                  "backgroundOpacity": 0,
                  "displayAsButton": False,
                  "buttonSize": "big",
                  "goToSlideType": "specified",
                  "invisible": False
                }
                    #h5p shape horizontal-line - theo chieu ngang
                    h5p_horizontal_line={
                  "x": 42.955892034233045,
                  "y": 49.84557214706733,
                  "width": 14.09,
                  "height": 27.831977000000002,
                  "action": {
                    "library": "H5P.Shape 1.0",
                    "params": {
                      "line": {
                        "borderWidth": 1,
                        "borderStyle": "solid",
                        "borderColor": "#000"
                      },
                      "type": "horizontal-line",
                      "shape": {
                        "fillColor": "#ffffff",
                        "borderColor": "#000000",
                        "borderWidth": "0",
                        "borderStyle": "solid",
                        "borderRadius": "0"
                      }
                    },
                    "subContentId": "fcb3f795-18b6-4a70-b864-2152e3026dee",
                    "metadata": {
                      "contentType": "Shapes"
                    }
                  },
                  "backgroundOpacity": 0,
                  "displayAsButton": False,
                  "buttonSize": "big",
                  "goToSlideType": "specified",
                  "invisible": False
                }

                    # h5p shape vertical-line - theo chieu doc
                    h5p_vertical_line={
                  "x": 18.104015799868332,
                  "y": 29.260014239873595,
                  "width": 14.09,
                  "height": 27.831977000000002,
                  "action": {
                    "library": "H5P.Shape 1.0",
                    "params": {
                      "line": {
                        "borderWidth": 1,
                        "borderStyle": "solid",
                        "borderColor": "#000"
                      },
                      "type": "vertical-line",
                      "shape": {
                        "fillColor": "#ffffff",
                        "borderColor": "#000000",
                        "borderWidth": "0",
                        "borderStyle": "solid",
                        "borderRadius": "0"
                      }
                    },
                    "subContentId": "619c1db3-4dd1-403b-9894-b1314541d96d",
                    "metadata": {
                      "contentType": "Shapes"
                    }
                  },
                  "backgroundOpacity": 0,
                  "displayAsButton": False,
                  "buttonSize": "big",
                  "goToSlideType": "specified",
                  "invisible": False
                }
                    #h5ptext
                    h5ptext = {
                                    "x": z.x,
                                    "y": z.y,
                                    "width": z.cx,
                                    "height": z.cy,
                                    "action": {
                                      "library": "H5P.AdvancedText 1.1",
                                      "params": {
                                        "text": cpmverttext_h5p(z.text,z.color,z.size)
                                        },
                                      "subContentId": str(nhap4.uuid.uuid4()),
                                      "metadata": {
                                        "contentType": "Text",
                                        "license": "U",
                                        "title": "Untitled Text",
                                        "authors": [],
                                        "changes": []
                                      }
                                    },
                                    "alwaysDisplayComments": False,
                                    "backgroundOpacity": 0,
                                    "displayAsButton": False,
                                    "buttonSize": "big",
                                    "goToSlideType": "specified",
                                    "invisible": False,
                                    "solution": ""
                                  } #% (z.x,z.y,z.cx,z.cy,nhap4.cpmverttext_h5p(z.text),str(nhap4.uuid.uuid4()))
                    print(h5ptext)
                    listobj_textfield.append(h5ptext)
                # uidramdom=str(nhap4.uuid.uuid4())
                # h5p_textt={"x":z.x,"y":z.y,"width":z.cx,"height":z.cy,"action":{"library":"H5P.AdvancedText 1.1","params":{"text":,"subContentId":uidramdom,"metadata":{"contentType":"Text","license":"U","title":"Untitled Text","authors":[],"changes":[]}},"alwaysDisplayComments":False,"backgroundOpacity":0,"displayAsButton":False,"buttonSize":"big","goToSlideType":"specified","invisible":False,"solution":""} #% (z.x,z.y,z.cx,z.cy,,str(uuid.uuid4()))
                #print('h5p_text_element',j,h5p_textt)



         #edit ở đây, phần lấy all slide ngày 23-9=2022
        listobj_text_slide.append(listobj_shapefield)




            #listobj_textfield = []
        print('#-------------END di thang obj-------------------')
                    #print(j,'z.text',z.text)
    # for i,x in enumerate(listobj_text_slide):
    #     if not x:
    #         print('slide ',i,'None')
    #      else:
    #         print('slide',i,' check slide after add')
    #         print('gia tri x',x)
    #         for j,data in enumerate(x):
    #             print('obj_text_h5p in di thang obj shape',data)
    #return listobj_text_slide#,list_none
    return listobj_text_slide#,listobj_shapefield#,list_none


def dithangobject_ver001(arrayobj,arr_none):
    #print('#-----di thang obj ver001-----------')
    #r"""\%s""" % ("")
    #x,y,width,height,text,subContentId
    #x,y,width,height= number -> nothing
    #text,subcontent is string -> ""
    print('#-------------di thang obj-------------------')
    listobj_textfield=[]
    listobj_text_slide=[]
    list_none=arr_none
    for i,x in enumerate(arrayobj):

        if not x:
            print('slide ',i,'None')
            listobj_text_slide.append('')
        else:
            print('slide truoc add', i)
            for j,z in enumerate(x):
                #cpmverttext_h5p
                print('check type',type(z.x))
                z.x, z.y, z.cx, z.cy = convert_donvi_text_xycxcy_2_h5p(z.x, z.y, z.cx, z.cy)

                h5ptext = {
                                "x": z.x,
                                "y": z.y,
                                "width": z.cx,
                                "height": z.cy,
                                "action": {
                                  "library": "H5P.AdvancedText 1.1",
                                  "params": {
                                    "text": cpmverttext_h5p(z.text,z.color,z.size)
                                    },
                                  "subContentId": str(nhap4.uuid.uuid4()),
                                  "metadata": {
                                    "contentType": "Text",
                                    "license": "U",
                                    "title": "Untitled Text",
                                    "authors": [],
                                    "changes": []
                                  }
                                },
                                "alwaysDisplayComments": False,
                                "backgroundOpacity": 0,
                                "displayAsButton": False,
                                "buttonSize": "big",
                                "goToSlideType": "specified",
                                "invisible": False,
                                "solution": ""
                              } #% (z.x,z.y,z.cx,z.cy,nhap4.cpmverttext_h5p(z.text),str(nhap4.uuid.uuid4()))
                print(h5ptext)
                listobj_textfield.append(h5ptext)
                # uidramdom=str(nhap4.uuid.uuid4())
                # h5p_textt={"x":z.x,"y":z.y,"width":z.cx,"height":z.cy,"action":{"library":"H5P.AdvancedText 1.1","params":{"text":,"subContentId":uidramdom,"metadata":{"contentType":"Text","license":"U","title":"Untitled Text","authors":[],"changes":[]}},"alwaysDisplayComments":False,"backgroundOpacity":0,"displayAsButton":False,"buttonSize":"big","goToSlideType":"specified","invisible":False,"solution":""} #% (z.x,z.y,z.cx,z.cy,,str(uuid.uuid4()))
                #print('h5p_text_element',j,h5p_textt)


            listobj_text_slide.append(listobj_textfield)
            listobj_textfield = []
            print('#-------------END di thang obj-------------------')
                #print(j,'z.text',z.text)
    for i,x in enumerate(listobj_text_slide):
        if not x:
            print('slide ',i,'None')
        else:
            print('slide',i)
            for j,data in enumerate(x):
                print(j,'obj_text_h5p in di thang obj',data)
    return listobj_text_slide#,list_none
def count_keys(dict):
    count = 0
    for i in enumerate(dict):
        count += 1
    return count
def dithangobject(arrayobj,arr_none):
    #r"""\%s""" % ("")
    #x,y,width,height,text,subContentId
    #x,y,width,height= number -> nothing
    #text,subcontent is string -> ""

    listobj_textfield=[]
    listobj_text_slide=[]
    list_none=arr_none
    for i,x in enumerate(arrayobj):

        if not x:
            print('slide ',i,'None')
        else:
            print('slide truoc add', i)
            for j,z in enumerate(x):
                #cpmverttext_h5p
                h5ptext = r"""{
                                "x": %s,
                                "y": %s,
                                "width": %s,
                                "height": %s,
                                "action": {
                                  "library": "H5P.AdvancedText 1.1",
                                  "params": {
                                    "text": "%s"      
                                    },
                                  "subContentId": "%s",
                                  "metadata": {
                                    "contentType": "Text",
                                    "license": "U",
                                    "title": "Untitled Text",
                                    "authors": [],
                                    "changes": []
                                  }
                                },
                                "alwaysDisplayComments": false,
                                "backgroundOpacity": 0,
                                "displayAsButton": false,
                                "buttonSize": "big",
                                "goToSlideType": "specified",
                                "invisible": false,
                                "solution": ""
                              }""" % (z.x,z.y,z.cx,z.cy,cpmverttext_h5p(z.text),str(uuid.uuid4()))
                h5p_textt="""{"x":%s,"y":%s,"width":%s,"height":%s,"action":{"library":"H5P.AdvancedText 1.1","params":{"text":"%s"},"subContentId":"%s","metadata":{"contentType":"Text","license":"U","title":"Untitled Text","authors":[],"changes":[]}},"alwaysDisplayComments":false,"backgroundOpacity":0,"displayAsButton":false,"buttonSize":"big","goToSlideType":"specified","invisible":false,"solution":""}""" % (z.x,z.y,z.cx,z.cy,cpmverttext_h5p(z.text),str(uuid.uuid4()))
                #print('h5p_text_element',j,h5p_textt)
                listobj_textfield.append(h5p_textt)

        listobj_text_slide.append(listobj_textfield)
        listobj_textfield = []
            #print(j,'z.text',z.text)
    # for i,x in enumerate(listobj_text_slide):
    #     if not x:
    #         print('slide ',i,'None')
    #     else:
    #         print('slide',i)
    #         for j,data in enumerate(x):
    #             print(j,'obj_text_h5p',data)
    return listobj_text_slide#,list_none
def ppt2png(file):
    try:
        powerpoint = client.Dispatch("Powerpoint.Application")
    except Exception as e:
        print("Powerpoint could not be opened", file=sys.stderr)
        raise e
    try: # look if an active presentation is open
        assert powerpoint.ActivePresentation is not None
        QUIT = False # don't quit powerpoint later
    except:
        QUIT = True # quit powerpoint later
    ppt = powerpoint.Presentations.Open(file)
    ppt.Export(file, "PNG")
    ppt.Close()
    if QUIT: # quit only if required
        powerpoint.Quit()
####lưu tâm mốt sửa sau!!!! ôi buồn quá!
# def edit_name_pptx2rar(title,folder):
#     original = r'C:\Users\Ron\Desktop\Test_1\products.csv'
#     target = r'C:\Users\Ron\Desktop\Test_2\products.csv'
#     shutil.copyfile(original, target)
#     old_name = str(title) + r"C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p-1.2_dêp\a - Copy.pptx"
#     new_name = str(folder) + r"C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p-1.2_dêp\a - Copy.rar"
#     # Renaming the file
#     os.rename(old_name, new_name)

# method tự động sinh text từ pptx to h5p
def converttext_h5p_mini(text,color):
    return r"""<span><span><span><span><span style=\"color:#%s\"><span><span><span>%s</span></span></span></span> """%(str(color),text)
def replacestring(text):

    text = text.replace(r"\\\%s"%(''), r"\%s"%(''))

    return text
def cpmverttext_h5p(text,color,size):
    #read xml -> add obj to array_obj
    #load_xml_from_folder()
    html=''

    a = converttext_h5p_mini(text,color)
    html = html+str(a)

    return r"""<p style=\"text-align:left\"><span style=\"font-size:%s\">%s</span></p>"""%('1em',html)
def xx(content,array_h5p,slide):
    aa = [i for i, x in enumerate(array_h5p) if not x]

    for numb, (a, b) in enumerate(zip(array_h5p, content["presentation"]["slides"])):
        if numb not in aa:
            if not a:
                print('slide',numb,'none')
            else:
                print('slide', numb)
                for j in a:

                    b["elements"].append(j)
    f = open("myfilee.txt", "x")

    f.close()
    with open('myfilee.txt', 'w', encoding='utf-8') as f:
        f.write(json.dumps(content))
    suoc_3 = r'\%s'%(r"""\\""")
    suoc_1 = r'\%s'%(r"""""")

    #print("The ASCII value of '" + c + "' is", ord(c))
    with open('myfilee.txt', 'r',encoding='utf-8') as file:
        filedata = file.read()

    # Replace the target string
    filedata = filedata.replace(suoc_3, suoc_1)

    # Write the file out again
    with open('myfilee.txt', 'w',encoding='utf-8') as file:
        file.write(filedata)
    os.rename('myfilee.txt', 'myfilee.json')
    with open('myfilee.json', encoding='utf-8') as fp:
        data = json.load(fp)

    os.remove("myfilee.json")
    return data
def add_rong_elements_0(list):
    for i, dataa in enumerate(list):
        if not i:
            print("k ton tai")
        else:
            dataa.insert(0, '')
            print('slide i', i)
            for numb, j in enumerate(dataa):
                # j.insert(0,'')
                if not j:
                    print('obj', numb, ' is none')
                else:
                    print(numb, j)

    return list
def traveobj_slide(title,folder):
    #folder = r'C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p_1_2_deep'
    #title = '10slide'

    newname_fodler_rar, folderr, new_title_rar, ten_folder_rar = edit_name_pptx2rar(title, folder)
    print('newname_fodler_rar:', newname_fodler_rar)
    print('new_title_rar:', new_title_rar)
    print('folderr:', folderr)
    print('ten_folder_rar:', ten_folder_rar)
    # --------------------------------------
    # xử lý lại extra file sau
    # extra_rarfile(new_title_rar, folderr)
    # --------------------------------------
    print('extra rar file')
    # load_xml_from_folder(os.path.splitext(filepath)[0])
    # #load link xml
    name_XMLa, linkfile_xml, xxxxx = load_xml_from_folder(
        os.path.join(folder, title))
    # link_images_pptx = load_images_from_folder_pptx_best(
    #     os.path.join(folder, title))
    #copy folder images from pptx add zip_h5p.
    #os.listdir(folder + str("""_copy\ppt""") + str("""\media"""))
    #print(os.listdir(folder +str( """_copy\ppt""") + str("""\slides""")))
    #END copy images to zip h5p from pptx media

    print("loadxml from folder", name_XMLa)
    print("link file xml", linkfile_xml)
    print("gia tri:", linkfile_xml[1])
    print('#------------------------------------')
    print("xxxxx", xxxxx)
    print('#------------------------------------')
    listfiletext = []
    list_file_obj_shape_imgg=[]
#--------------- sửa lại ở đây - 30-8
    for i in xxxxx:#lay_all_objtext_infilexml
        #listfiletext.append(elementree_withedittextt.lay_all_objtext_infilexml(os.path.join(folder, title), i))
        #30-8 đã sửa lại từ lay__all_obj_inxml từ text tới text+ shape . . .
        #lay all file _obj text của elementree_width-text này có cả phần lấy color của chữ, sau xem lại!
        # nó đây .elementree_withedittextt.lay_all_objtext_infilexml
        print('xu ly xxxxx')
        objtext,objshape=elementree_withedittextt.lay_all_objtext_infilexml(os.path.join(folder, title), i)
        listfiletext.append(objtext)
        list_file_obj_shape_imgg.append(objshape)

    a = [i for i, x in enumerate(listfiletext) if not x]
    print('obj None', a)
    # for i,x in enumerate(listfiletext):
    #
    #     if not x:
    #         print('slide',i,'none')
    #     else:
    #         print('slidei', i)
    #         for j,z in enumerate(x):
    #             print('element j',j,'z.text',z.text)


    #KHONG LAY STRING OBJECT NƯA
    # aaa = dithangobject(listfiletext, a)
    # for i, x in enumerate(aaa):
    #     if not x:
    #         print('slide ', i, 'None')
    #     else:
    #         print('slide', i)
    #         for j, data in enumerate(x):
    #             print(j, 'obj_text_h5p', data,'\n')

    return listfiletext,list_file_obj_shape_imgg
# def add_to_json(newfile, image_folder, images, title, arrayobj):
#     exclude_files = ['content/content.json', 'h5p.json']
#     if getattr(sys, 'frozen', False):
#         basedir = sys._MEIPASS
#     else:
#         basedir = os.path.dirname(os.path.abspath(__file__))
#     template = os.path.join(basedir, "template.h5p")
#
#     img_width, img_height = get_image_size(os.path.join(image_folder, images[0]))
#
#     with ZipFile(template, "r") as zin:
#         with ZipFile(newfile, "w") as zout:
#             # copy all other files
#             for item in zin.infolist():
#                 if item.filename not in exclude_files:
#                     zout.writestr(item, zin.read(item.filename))
#
#             # add image filenames to content.json
#             with zin.open("content/content.json") as fp:
#                 content = json.load(fp)
#                 #----------------------
#                 #print('value content is:',content)
#                 #print('value content dump:',json.dumps(content, indent=1))
#                 # ----------------------
#                 # for ii in content:
#                 #     print('value ii',ii)
#                 #print('conten',content)
#                 print(f"adding images from {image_folder}: {images}")
#                 for i, image in enumerate(images):
#                     slides = content["presentation"]["slides"]
#                     # clone first element
#                     if i > 0 and len(slides) <= i:
#                         slides.append(deepcopy(slides[0]))
#
#                     # print(i,'value file images slide',slides[i]["elements"][0])
#                     # print('#-------------------------------------------imgaes---')
#                     #
#                     # print(i,'value file images slide',slides[i]["elements"][1])
#                     elements = slides[i]["elements"][0]
#                     params = elements["action"]["params"]
#                     # add new image filename
#                     params["file"]["path"] = "images/" + image
#                     # add random uuid
#                     elements["action"]["subContentId"] = str(uuid.uuid4())
#                     #set width & height
#                     params["file"]["width"] = img_width
#                     params["file"]["height"] = img_height
#                     ratio = img_width / img_height
#                     if ratio > target_ratio:  # wider, need to shrink y
#                         elements["y"] = 100 * (1 - target_ratio / ratio) / 2
#                         elements["height"] = 100 * target_ratio / ratio
#                     elif ratio < target_ratio:  # higher, need to shrink x
#                         elements["x"] = 100 * (1 - ratio / target_ratio) / 2
#                         elements["width"] = 100 * ratio / target_ratio
#                 # #god bles you ♥
#                             # for i, x in enumerate(listt):
#                             #
#                             #     print(i, x)
#                             #     print('#-------------------------text obj.x.y-------------------')
#                             #
#                             #     for j in x:
#                             #         print(i, j, ',', j.text, j.x, j.y)
#                             #         # for z in j:
#                             #         #  print('x:',z.text)
#                             #
#                             #     print('#-------------------------end text obj.x.y-------------------')
#                 #slides = content["presentation"]["slides"]
#
#                 # for x in slides:
#                 #     print('x in slides',x)
#                 # print('print slide 0:',slides[0])
#                 # print('print slide 1:',slides[1])
#
#                 # #god bles you ♥
#                 # sodem=0
#                 # slides = content["presentation"]["slides"]  # ["elements"]
#                 # for (i_obj, slide) in itertools.zip_longest(arrayobj, slides):
#                 #     if sodem in a:
#                 #         continue
#                 #
#                 #     for x in enumerate(i_obj):
#                 #         slide['elements'].append(x)
#
#
#                     #print(i, j)
#                 a = [i for i, x in enumerate(arrayobj) if not x]
#                 #print('value of slide',slides.count('elements'))
#                 for x in a:
#                     print('values arrayobj none:', x)
#                 for i, (a, b) in enumerate(zip(arrayobj, content["presentation"]["slides"])):
#                     print(i, 'value cua b',b, )
#                     b['elements'].append(a)
#                     print(b['elements'],'new B elemtns')
#                 # for i, (arr_obj, slide) in enumerate(zip(arrayobj, content["presentation"]["slides"])):
#                 #
#                 #     if i in a:
#                 #         continue
#                 #     else:
#                 #         print(i)
#                 #         for j, arr_obj_elements in enumerate(arr_obj):
#                 #
#                 #             #print(j,'text-cx-cy',dataa.text,dataa.cx,dataa.cy)
#                 #             print(j,'data obj apen in slide ',i," la: ",arr_obj_elements)
#                 #             slide.append(arr_obj_elements)
#
#                 # for i, arr_obj in enumerate(arrayobj,slides):
#                 #     if i in a:
#                 #         continue
#                 #     else:
#                 #         print(i)
#                 #         for j,dataa in enumerate(arr_obj):
#                 #
#                 #             #print(j,'text-cx-cy',dataa.text,dataa.cx,dataa.cy)
#                 #             print(j,'data obj',dataa)
#                         #xml after cv 2 json- lấy hết obj from từng xml
#                         #obj_xml_cv2json,obj_rutgon = xml2json.xml_2_json_getfromfile3(link_xml_file)
#                 #print('slide', i)
#
#                 print('#-------------------------text obj.x.y-------------------')
#                         #get element to add param text to json
#                         #edit element slides ♥ god bles u
# #                       #slides = content["presentation"]["slides"]
#                         # clone first element
#                         # for x in slides:
#                         #     print('x in slides',x)
#                         #print()
#                         #tính lại x,y, cx,cy
#                         #viet method
#                         #convert_donvi_text_xycxcy_2_h5p()
#                 print('#-------------------------end text obj.x.y-------------------')
#
#
#                         #if not arr_obj:
#                         #print(i, 'i none')
#                         #continue
#                         #else:
#
#
#                         # for j,data in enumerate(arr_obj) or []:
#                         #     #add element text cho tung slide
#                         #     #print("gia tri cua value element:", j, data)
#                         #     if j > 0 and len(arr_obj) <= j :
#                         #         #xx = xx + 1
#                         #         slides[i]["elements"].append(deepcopy(slides[0]["element"][1]))
#                                 #slides[i]['elements'].append(data)
#                 #print(r"slides[4][elements][2]",slides[4]["elements"][2])
#                 #slides = content["presentation"]["slides"]
#                 # for i, arr_obj in enumerate(slides):
#                 #
#                 #     for j , daa in enumerate(arr_obj):
#                 #         print('in ra value slide:')
#                 #         print('slide',i,'element ',j,'gia tri cua value element sau apend', daa[i]['elements'][j])
#
#                     #tính lại x,y, cx,cy
#                             # viet method
#                     # if(1==1):
#                     #     #----------------------------------------- phai xu ly duoc list index
#                     #     data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
#                     #     elements = slides[i]["elements"][j+1]
#                     #     slides[i]["elements"][j+1]['x']=data.x#['x']
#                     #     slides[i]["elements"][j+1]['y'] = data.y#['y']
#                     #     slides[i]["elements"][j+1]['width'] = data.cx#['cx']
#                     #     slides[i]["elements"][j+1]['height'] = data.cy#['cy']
#                     #
#                     #     params = elements["action"]["params"]
#                     #     # add new image filename
#                     #     #---------------------------
#                     #     params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
#                     #     #---------------------------
#                     #     # add random uuid
#                     #     elements["action"]["subContentId"] = str(uuid.uuid4())
#                         # set width & height
#                         # params["file"]["width"] = img_width
#                         # params["file"]["height"] = img_height
#                 print('#-------------------------end text obj.x.y-------------------')
#                                 #else:
#                                 #xx=xx-1
#                             #xx=xx+1
#                         #
#
#                 # for i, arr_obj in enumerate(arrayobj):
#                 #     for j, data in enumerate(arr_obj):
#                 #
#                 #         # tính lại x,y, cx,cy
#                 #         # viet method
#                 #         if(1==1 ):
#
#                             # data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
#                             # elements = slides[i]["elements"][j]
#                             # slides[i]["elements"][j]['x']=data.x#['x']
#                             # slides[i]["elements"][j]['y'] = data.y#['y']
#                             # slides[i]["elements"][j]['width'] = data.cx#['cx']
#                             # slides[i]["elements"][j]['height'] = data.cy#['cy']
#                             #
#                             # params = elements["action"]["params"]
#                             # # add new image filename
#                             # #---------------------------
#                             # params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
#                             # #---------------------------
#                             # # add random uuid
#                             # elements["action"]["subContentId"] = str(uuid.uuid4())
#                             # set width & height
#                             # params["file"]["width"] = img_width
#                             # params["file"]["height"] = img_height
#                 # print('#-------------------------end text obj.x.y-------------------')
#                 for i, image in enumerate(images):
#                     slides = content["presentation"]["slides"]
#                     # clone first element
#                     # if i > 0 and len(slides) <= i:
#                     #     slides.append(deepcopy(slides[0]))
#                     elements = slides[i]["elements"][0]
#                     params = elements["action"]["params"]
#                     # add new image filename
#                     params["file"]["path"] = "images/" + image
#                     # add random uuid
#                     elements["action"]["subContentId"] = str(uuid.uuid4())
#                     #set width & height
#                     params["file"]["width"] = img_width
#                     params["file"]["height"] = img_height
#                     ratio = img_width / img_height
#                     if ratio > target_ratio:  # wider, need to shrink y
#                         elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
#                         elements["height"] = 100# * target_ratio / ratio
#                     elif ratio < target_ratio:  # higher, need to shrink x
#                         elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
#                         elements["width"] = 100 #* ratio / target_ratio
#                     #print(i,'value file images slide',slides[i]["elements"][0])
#             # save file
#             zout.writestr("content/content.json", json.dumps(content))
#
#             # add image files to tip
#             for image in images:
#                 zout.write(os.path.join(image_folder, image), "content/images/" + image)
#
#             # change presentation title
#             with zin.open("h5p.json", "r") as h5p:
#                 content = json.load(h5p)
#                 content["title"] = title
#             zout.writestr("h5p.json", json.dumps(content))
#
#
#
#             # save file
#             zout.writestr("content/content.json", json.dumps(content))
#
#             # add image files to tip
#             for image in images:
#                 zout.write(os.path.join(image_folder, image), "content/images/" + image)
#
#             # change presentation title
#             with zin.open("h5p.json", "r") as h5p:
#                 content = json.load(h5p)
#                 content["title"] = title
#             zout.writestr("h5p.json", json.dumps(content))
def convert_donvi_text_xycxcy_2_h5p(x,y,cx,cy):
    # p: sldSz
    # type = "screen4x3"
    # cy = "6858000"
    # cx = "9144000" /

#    print (type (x))
    # xx=(100*int(x))/9144000
    # yy=(100*int(y))/6858000
    # cxx = (100 * int(cx)) / 9144000
    # cyy = (100 * int(cy)) / 6858000
    print('new value off donvitextt x y cx cy=',x,y,cx,cy)
    # x=str((100*float(x))/9144000)
    # y=str((100*float(y))/6858000)
    # cx = str((100 * float(cx)) / 9144000)
    # cy = str((100 * float(cy)) / 6858000)
    # x = str((100.0 * int(x)) / 9144000.0)
    #
    # y = str((100.0 * int(y)) / 6858000.0)
    # cx = str((100.0 * int(cx)) / 9144000.0)
    # cy = str((100.0 * int(cy)) / 6858000.0)
    # if xx<0:
    #     xx=0
    # if yy < 0:
    #     yy = 0
    # if cxx < 0:
    #     cxx = 0
    # if cyy < 0:
    #     cyy = 0

    if x[0] == '-':
        x = str(0)
    if y[0] == '-':
        y = str(0)
    if cx [0] == '-':
        cx = str(0)
    if cy[0] == '-':
        cy = str(0)

    return x,y,cx,cy
    #return abs(xx),abs(yy),abs(cxx),abs(cyy)
def khoitao_element_rong(arr):
    for i in arr:
        for j in i:
            j.insert(0, '')
    return arr
def dequyaddelements(content,elements):


    return content


def edit_elements_slide(data, x):
    a = [i for i, x in enumerate(x) if not x]
    print(a)
    for i, (slide, arr_obj) in enumerate(zip(data["presentation"]["slides"], x)):
        print('#------------------------start with h5p add text')
        if (True == True):
            if i not in a:
                print('slide', i)
                # print("gia tri cua arrayobj",i,arr_obj)
                # print('#-------------------------text obj.x.y-------------------')

                # get element to add param text to json
                # edit element slides ♥ god bles u
                # tính lại x,y, cx,cy
                # viet method
                # convert_donvi_text_xycxcy_2_h5p()
                # print('#-------------------------end text obj.x.y-------------------')
                # for j,data in enumerate(arr_obj):

                for j, (elements, data) in enumerate(zip(slide['elements'], arr_obj)):
                    # add element text cho tung slide

                    # tính lại x,y, cx,cy
                    #         viet method
                    if (True):
                        if data == '':
                            print('data', j, ' is none\n')
                        else:

                            # ----------------------------------------- phai xu ly duoc list index
                            print('#-----------start convert x y cx cy---------------')
                            print('gia tri array[', i, '],[', j, ']')

                            print(data.x, data.y, data.cx, data.cy, data.text)
                            data.x, data.y, data.cx, data.cy = convert_donvi_text_xycxcy_2_h5p(data.x, data.y, data.cx,
                                                                                               data.cy)
                            print(data.x, data.y, data.cx, data.cy, data.text)
                            print('slide[', i, ']', 'elements[', j, ']', elements)

                            if (False):
                                print('gia tri cua elements object')
                                print(elements['x'],  # ['x']
                                      elements['y'],  # ['y']
                                      elements['width'],  # ['cx']
                                      elements['height']  # ['cy'])
                                      ,
                                      elements["action"]["params"]
                                      )
                                params = elements["action"]["params"]
                                params.update({"text": cpmverttext_h5p(data.text)})
                                print('new value parm text', params)
                                print('value elements["action"]["subContentId"] ', elements["action"]["subContentId"])
                                elements["action"]["subContentId"] = str(uuid.uuid4())
                                print('value elements["action"]["subContentId"] ', elements["action"]["subContentId"])
                            print('#-----------end convert x y cx cy-----------------')
                            # params = elements["action"]["params"]
                            # params.update({"text": replacestring(cpmverttext_h5p(data.text))})
                            # print('text',params)
                            # slides[0]["elements"][1]
                            # elements = slides[i]["elements"][j+1]
                            # checkslide_arr(elements,data)
                            if (True):
                                elements['x'] = data.x  # ['x']
                                elements['y'] = data.y  # ['y']
                                elements['width'] = data.cx  # ['cx']
                                elements['height'] = data.cy  # ['cy']

                                params = elements["action"]["params"]
                                # add new image filename
                                # ---------------------------
                                # params["text"] = replacestring(cpmverttext_h5p(cpmverttext_h5p(data.text)))
                                params.update({"text": cpmverttext_h5p(data.text)})
                                # ---------------------------
                                # add random uuid
                                elements["action"]["subContentId"] = str(uuid.uuid4())
    return data
def add_to_newjson(newfile, image_folder, images, title, arrayobj):

    exclude_files = ['content/content.json', 'h5p.json']
    if getattr(sys, 'frozen', False):
        basedir = sys._MEIPASS
    else:
        basedir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(basedir, newfile)

    img_width, img_height = get_image_size(os.path.join(image_folder, images[0]))

    with ZipFile(template, "r") as zin:
        with ZipFile(newfile, "w") as zout:
            # copy all other files
            for item in zin.infolist():
                print("item.filename",item.filename)
                if item.filename not in exclude_files:
                    print(zin.read(item.filename))
                    zout.writestr(item, zin.read(item.filename))

            # add image filenames to content.json
            with zin.open("content/content.json") as fp:
                content = json.load(fp)
                #----------------------
                #print('value content is:',content)
                #print('value content dump:',json.dumps(content, indent=1))
                # ----------------------
                # for ii in content:
                #     print('value ii',ii)
                #print('conten',content)
                print(f"adding images from {image_folder}: {images}")
                for i, image in enumerate(images):
                    slides = content["presentation"]["slides"]
                    # clone first element
                    if i > 0 and len(slides) <= i:
                        slides.append(deepcopy(slides[0]))

                    # print(i,'value file images slide',slides[i]["elements"][0])
                    # print('#-------------------------------------------imgaes---')
                    #
                    # print(i,'value file images slide',slides[i]["elements"][1])
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100 * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 * ratio / target_ratio
                # #god bles you ♥
                            # for i, x in enumerate(listt):
                            #
                            #     print(i, x)
                            #     print('#-------------------------text obj.x.y-------------------')
                            #
                            #     for j in x:
                            #         print(i, j, ',', j.text, j.x, j.y)
                            #         # for z in j:
                            #         #  print('x:',z.text)
                            #
                            #     print('#-------------------------end text obj.x.y-------------------')
                #slides = content["presentation"]["slides"]

                # for x in slides:
                #     print('x in slides',x)
                # print('print slide 0:',slides[0])
                # print('print slide 1:',slides[1])
                # #god bles you ♥
                slides = content["presentation"]["slides"]  # ["elements"]

                # a = [i for i, x in enumerate(arrayobj) if not x]
                # for x in a:
                #     print('values arrayobj none:', x)
                # for i, arr_obj in enumerate(arrayobj):
                #     if i in a:
                #         continue
                #     else:
                        # print(i)
                        # for j,dataa in enumerate(arr_obj):
                        #     print(j,'text-cx-cy',dataa.text,dataa.cx,dataa.cy)
                        #xml after cv 2 json- lấy hết obj from từng xml
                        #obj_xml_cv2json,obj_rutgon = xml2json.xml_2_json_getfromfile3(link_xml_file)
                #print('slide', i)

                print('#-------------------------text obj.x.y-------------------')
                        #get element to add param text to json
                        #edit element slides ♥ god bles u
#                       #slides = content["presentation"]["slides"]
                        # clone first element
                        # for x in slides:
                        #     print('x in slides',x)
                        #print()
                        #tính lại x,y, cx,cy
                        #viet method
                        #convert_donvi_text_xycxcy_2_h5p()
                print('#-------------------------end text obj.x.y-------------------')


                        #if not arr_obj:
                        #print(i, 'i none')
                        #continue
                        #else:


                        # for j,data in enumerate(arr_obj) or []:
                        #     #add element text cho tung slide
                        #     #print("gia tri cua value element:", j, data)
                        #     if j > 0 and len(arr_obj) <= j :
                        #         #xx = xx + 1
                        #         slides[i]["elements"].append(deepcopy(slides[0]["element"][1]))
                                #slides[i]['elements'].append(data)
                #print(r"slides[4][elements][2]",slides[4]["elements"][2])
                #slides = content["presentation"]["slides"]
                # for i, arr_obj in enumerate(slides):
                #
                #     for j , daa in enumerate(arr_obj):
                #         print('in ra value slide:')
                #         print('slide',i,'element ',j,'gia tri cua value element sau apend', daa[i]['elements'][j])

                    #tính lại x,y, cx,cy
                            # viet method
                    # if(1==1):
                    #     #----------------------------------------- phai xu ly duoc list index
                    #     data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                    #     elements = slides[i]["elements"][j+1]
                    #     slides[i]["elements"][j+1]['x']=data.x#['x']
                    #     slides[i]["elements"][j+1]['y'] = data.y#['y']
                    #     slides[i]["elements"][j+1]['width'] = data.cx#['cx']
                    #     slides[i]["elements"][j+1]['height'] = data.cy#['cy']
                    #
                    #     params = elements["action"]["params"]
                    #     # add new image filename
                    #     #---------------------------
                    #     params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
                    #     #---------------------------
                    #     # add random uuid
                    #     elements["action"]["subContentId"] = str(uuid.uuid4())
                        # set width & height
                        # params["file"]["width"] = img_width
                        # params["file"]["height"] = img_height
                print('#-------------------------end text obj.x.y-------------------')
                                #else:
                                #xx=xx-1
                            #xx=xx+1
                        #
                for numb, i in enumerate(slides["presentation"]["slides"]):
                    # print('slide',numb)
                    print('value numb', numb)
                    for numj, j in enumerate(i['elements']):
                        print(numj, j)
                # for i, arr_obj in enumerate(arrayobj,["presentation"]["slides"]):
                #     for j, data in enumerate(arr_obj):
                #
                #         # tính lại x,y, cx,cy
                #         # viet method
                #         if(1==1 ):

                            # data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                            # elements = slides[i]["elements"][j]
                            # slides[i]["elements"][j]['x']=data.x#['x']
                            # slides[i]["elements"][j]['y'] = data.y#['y']
                            # slides[i]["elements"][j]['width'] = data.cx#['cx']
                            # slides[i]["elements"][j]['height'] = data.cy#['cy']
                            #
                            # params = elements["action"]["params"]
                            # add new image filename
                            #---------------------------
                            # params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
                            # #---------------------------
                            # # add random uuid
                            # elements["action"]["subContentId"] = str(uuid.uuid4())
                #             # set width & height
                #             # params["file"]["width"] = img_width
                #             # params["file"]["height"] = img_height
                # print('#-------------------------end text obj.x.y-------------------')
                for i, image in enumerate(images):
                    slides = content["presentation"]["slides"]
                    # clone first element
                    # if i > 0 and len(slides) <= i:
                    #     slides.append(deepcopy(slides[0]))
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100# * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 #* ratio / target_ratio
                    #print(i,'value file images slide',slides[i]["elements"][0])
            # save file
            zout.writestr("content/content.json", json.dumps(content))

            # add image files to tip
            for image in images:
                zout.write(os.path.join(image_folder, image), "content/images/" + image)

            # change presentation title
            with zin.open("h5p.json", "r") as h5p:
                content = json.load(h5p)
                content["title"] = title
            zout.writestr("h5p.json", json.dumps(content))
def checkslide_arr(element,arr):
    print('#-------------------check list slide va arr:')

    print('value of slide:',element)
    print('value of slide use slide.value:', element.value())

    print('value of elements:',element)
    for i in element:
        print(i.values())
    for i,data in enumerate(arr):
        print(i,'data',data)
    print('#-------------------end check list slide va arr:')

def add_to_json(newfile, image_folder, images, title, arrayobj):

    exclude_files = ['content/content.json', 'h5p.json']
    if getattr(sys, 'frozen', False):
        basedir = sys._MEIPASS
    else:
        basedir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(basedir, "template_new.h5p")

    img_width, img_height = get_image_size(os.path.join(image_folder, images[0]))

    with ZipFile(template, "r") as zin:
        with ZipFile(newfile, "w") as zout:
            # copy all other files
            for item in zin.infolist():
                print("item.filename",item.filename)
                if item.filename not in exclude_files:
                    print(zin.read(item.filename))
                    zout.writestr(item, zin.read(item.filename))

            # add image filenames to content.json
            with zin.open("content/content.json") as fp:
                content = json.load(fp)
                #----------------------
                print('value content is:',content)
                print('value content dump:',json.dumps(content, indent=1))
                # ----------------------
                # for ii in content:
                #     print('value ii',ii)
                #print('conten',content)
                slides = content["presentation"]["slides"]
                print(f"adding images from {image_folder}: {images}")
                for i, image in enumerate(images):
                    #slides = content["presentation"]["slides"]
                    # clone first element
                    if i > 0 and len(slides) <= i:
                        slides.append(deepcopy(slides[0]))
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100 * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 * ratio / target_ratio


                # #god bles you ♥
                a = [i for i, x in enumerate(arrayobj) if not x]
                for x in a:
                    print('values arrayobj none:', x)
                #add elements slide
                for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                    if i in a:
                        print('slide ',i,'is none')
                    else:
                        print('slide ',i)
                        #print('values slide 0 element 1',slides[0]["elements"][1])
                        for j,data in enumerate(arr_obj):
                            print('value j',j)
                            #and j <= count_keys(data)
                            if(j>=0):
                                slide["elements"].append(slides[0]["elements"][1])
                # delete elements neu can
                delete_text_element1_inslide(content)


                # for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                #     if i in a:
                #         print('slide ',i,'is none')
                #     else:
                #         print('slide ',i)
                #         for j,data in enumerate(arr_obj):
                #             slide["elements"].append(slides[0]["elements"][1])
                #     #del slides[i]["elements"][1]
                number=0
                for i in content["presentation"]["slides"]:
                    print('\nslide',number)
                    number+=1
                    print(i['elements'])


                print('#------------------------start with add elemtens--------')
                # list01 = [1, 2, 3]
                # list02 = ['x','y']
                # for i, (a, b) in enumerate(zip(list01, list02)):
                #     print (i, a, b)

                add_rong_elements_0(arrayobj)
                print('#--------------check  sau khi add rong arr i-----------')
                for i,data in enumerate(arrayobj):
                    print(i)
                    for j,dataa in enumerate(data):
                        if j>0:
                            print('#--------------arr j-----------')
                            #,dataa.y,dataa.cx,dataa.cy,dataa.text
                            print(dataa.x,dataa.y,dataa.cx,dataa.cy,dataa.text)
                            print('#--------------arr j-----------')
                print('#--------------check  sau khi add rong end arr i-----------')
                #edit_elements_slide(content,arrayobj)
                if(False):
                    for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                        print('#------------------------start with h5p add text')
                        if(True==True):
                            if i not in a:
                                print('slide', i)
                                #print("gia tri cua arrayobj",i,arr_obj)
                                #print('#-------------------------text obj.x.y-------------------')

                                #get element to add param text to json
                                #edit element slides ♥ god bles u
                                #tính lại x,y, cx,cy
                                #viet method
                                #convert_donvi_text_xycxcy_2_h5p()
                                #print('#-------------------------end text obj.x.y-------------------')
                                #for j,data in enumerate(arr_obj):

                                for j, (elements, data) in enumerate(zip(slide['elements'], arr_obj)):
                                    #add element text cho tung slide


                                    # tính lại x,y, cx,cy
                                    #         viet method
                                    if(True):
                                        if data=='':
                                            print('data',j,' is none\n')
                                        else:

                                            #----------------------------------------- phai xu ly duoc list index
                                            print('#-----------start convert x y cx cy---------------')
                                            print('gia tri array[',i,'],[',j,']')

                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            print('slide[',i,']','elements[',j,']',elements)
                                            print('#-----------end convert x y cx cy-----------------')

                                            #slides[0]["elements"][1]
                                            #elements = slides[i]["elements"][j+1]
                                            #checkslide_arr(elements,data)
                                            if(True):
                                                elements['x']=data.x#['x']
                                                elements['y'] = data.y#['y']
                                                elements['width'] = data.cx#['cx']
                                                elements['height'] = data.cy#['cy']

                                                params = elements["action"]["params"]
                                                # add new image filename
                                                #---------------------------
                                                #params["text"] = replacestring(cpmverttext_h5p(cpmverttext_h5p(data.text)))
                                                params.update({"text": cpmverttext_h5p(data.text)})
                                                #---------------------------
                                                # add random uuid
                                                elements["action"]["subContentId"] = str(uuid.uuid4())
                                            #set width & height
                                            # params["file"]["width"] = img_width
                                            # params["file"]["height"] = img_height
                                    print('#-------------------------end text obj.x.y-------------------')
                    print('#------------------------start with h5p add text')

                                #else:
                                #xx=xx-1
                            #xx=xx+1
                        #

                    # for i, arr_obj in enumerate(arrayobj):
                    #     for j, data in enumerate(arr_obj):
                    #
                    #         # tính lại x,y, cx,cy
                    #         # viet method
                    #         if(1==1 ):
                    #             data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                    #             elements = slides[i]["elements"][j]
                    #             slides[i]["elements"][j]['x']=data.x#['x']
                    #             slides[i]["elements"][j]['y'] = data.y#['y']
                    #             slides[i]["elements"][j]['width'] = data.cx#['cx']
                    #             slides[i]["elements"][j]['height'] = data.cy#['cy']
                    #
                    #             params = elements["action"]["params"]
                    #             # add new image filename
                    #             #---------------------------
                    #             params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
                    #             #---------------------------
                    #             # add random uuid
                    #             elements["action"]["subContentId"] = str(uuid.uuid4())
                    #             # set width & height
                    #             # params["file"]["width"] = img_width
                    #             # params["file"]["height"] = img_height
                    # print('#-------------------------end text obj.x.y-------------------')

                for i, image in enumerate(images):
                    slides = content["presentation"]["slides"]
                    # clone first element
                    # if i > 0 and len(slides) <= i:
                    #     slides.append(deepcopy(slides[0]))
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100# * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 #* ratio / target_ratio
            # save file
            #content =
            zout.writestr("content/content.json", json.dumps(content))

            # add image files to tip
            for image in images:
                zout.write(os.path.join(image_folder, image), "content/images/" + image)

            # change presentation title
            with zin.open("h5p.json", "r") as h5p:
                content = json.load(h5p)
                content["title"] = title
            zout.writestr("h5p.json", json.dumps(content))
# def returndict(arr):
#
#     return objectt
def replaceslash(original):
    with open(original, 'r', encoding='utf-8') as file:
        filedata = file.read()

    filedata = filedata.replace(r"""\\\%s""" % (''), r"""\%s""" % (''))
    filedata = filedata.replace(r"""\\%s""" % (''), r"""\%s""" % (''))

    with open(original, 'w', encoding='utf-8') as file:
        file.write(filedata)

def create_newjson_withoutslash(content):
    with open('readme.txt', 'w') as f:
        f.write(json.dumps(content))
    replaceslash('readme.txt')
    with io.open('readme.txt', 'r', encoding='utf8') as f:
        text = f.read()
    os.remove('readme.txt')
    return text
def add_to_json_newver1(newfile, image_folder, images, title, arrayobj):

    exclude_files = ['content/content.json', 'h5p.json']
    exclude_filess = ['content/content.json']

    if getattr(sys, 'frozen', False):
        basedir = sys._MEIPASS
    else:
        basedir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(basedir, "template_new.h5p")

    img_width, img_height = get_image_size(os.path.join(image_folder, images[0]))

    with ZipFile(template, "r") as zin:
        with ZipFile(newfile, "w") as zout:
            # copy all other files
            for item in zin.infolist():
                print("item.filename",item.filename)
                if item.filename not in exclude_files:
                    print(zin.read(item.filename))
                    zout.writestr(item, zin.read(item.filename))

            # add image filenames to content.json
            with zin.open("content/content.json") as fp:
                content = json.load(fp)

                #----------------------
                print('value content is:',content)
                print('value content dump:',json.dumps(content, indent=1))
                # ----------------------
                # for ii in content:
                #     print('value ii',ii)
                #print('conten',content)
                slides = content["presentation"]["slides"]
                print(f"adding images from {image_folder}: {images}")
                for i, image in enumerate(images):
                    #slides = content["presentation"]["slides"]
                    # clone first element
                    if i > 0 and len(slides) <= i:
                        slides.append(deepcopy(slides[0]))
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100# * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 #* ratio / target_ratio


                # #god bles you ♥
                a = [i for i, x in enumerate(arrayobj) if not x]
                for x in a:
                    print('values arrayobj none:', x)
                #add elements slide
                for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                    if i in a:
                        print('slide ',i,'is none')
                    else:
                        print('slide ',i)
                        #print('values slide 0 element 1',slides[0]["elements"][1])
                        for j,data in enumerate(arr_obj):
                            print('value j',j)
                            #and j <= count_keys(data)
                            #if(j>=0):
                            slide["elements"].append(data)
                delete_text_element1_inslide(content)
                # delete elements neu can
                #delete_text_element1_inslide(content)
                # for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                #     if i in a:
                #         print('slide ',i,'is none')
                #     else:
                #         print('slide ',i)
                #         for j,data in enumerate(arr_obj):
                #             slide["elements"].append(slides[0]["elements"][1])
                #     #del slides[i]["elements"][1]



                print('#------------------------start with add elemtens--------')
                # list01 = [1, 2, 3]
                # list02 = ['x','y']
                # for i, (a, b) in enumerate(zip(list01, list02)):
                #     print (i, a, b)

                #add_rong_elements_0(arrayobj)
                print('#-------------- start check   arr i-----------')
                for i,data in enumerate(arrayobj):
                    print('obj ',i)
                    for j,dataa in enumerate(data):

                        print('#--------------arr j-----------')
                        #,dataa.y,dataa.cx,dataa.cy,dataa.text
                        #print(dataa.x,dataa.y,dataa.cx,dataa.cy,dataa.text)
                        print('elemsnts',j,data)
                        print('#--------------arr j-----------')
                print('#--------------check  sau khi add rong end arr i-----------')
                #edit_elements_slide(content,arrayobj)

                if(False):
                    for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                        print('#------------------------start with h5p add text')
                        if(True==True):
                            if i not in a:
                                print('slide', i)
                                #print("gia tri cua arrayobj",i,arr_obj)
                                #print('#-------------------------text obj.x.y-------------------')

                                #get element to add param text to json
                                #edit element slides ♥ god bles u
                                #tính lại x,y, cx,cy
                                #viet method
                                #convert_donvi_text_xycxcy_2_h5p()
                                #print('#-------------------------end text obj.x.y-------------------')
                                #for j,data in enumerate(arr_obj):

                                for j, (elements, data) in enumerate(zip(slide['elements'], arr_obj)):
                                    #add element text cho tung slide


                                    # tính lại x,y, cx,cy
                                    #         viet method
                                    if(True):
                                        if data=='':
                                            print('data',j,' is none\n')
                                        else:

                                            #----------------------------------------- phai xu ly duoc list index
                                            print('#-----------start convert x y cx cy---------------')
                                            print('gia tri array[',i,'],[',j,']')

                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            print('slide[',i,']','elements[',j,']',elements)
                                            print('#-----------end convert x y cx cy-----------------')

                                            #slides[0]["elements"][1]
                                            #elements = slides[i]["elements"][j+1]
                                            #checkslide_arr(elements,data)
                                            if(True):
                                                elements['x']=data.x#['x']
                                                elements['y'] = data.y#['y']
                                                elements['width'] = data.cx#['cx']
                                                elements['height'] = data.cy#['cy']

                                                params = elements["action"]["params"]
                                                # add new image filename
                                                #---------------------------
                                                #params["text"] = replacestring(cpmverttext_h5p(cpmverttext_h5p(data.text)))
                                                params.update({"text": cpmverttext_h5p(data.text)})
                                                #---------------------------
                                                # add random uuid
                                                elements["action"]["subContentId"] = str(uuid.uuid4())
                                            #set width & height
                                            # params["file"]["width"] = img_width
                                            # params["file"]["height"] = img_height
                                    print('#-------------------------end text obj.x.y-------------------')
                    print('#------------------------start with h5p add text')

                                #else:
                                #xx=xx-1
                            #xx=xx+1
                        #

                    # for i, arr_obj in enumerate(arrayobj):
                    #     for j, data in enumerate(arr_obj):
                    #
                    #         # tính lại x,y, cx,cy
                    #         # viet method
                    #         if(1==1 ):
                    #             data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                    #             elements = slides[i]["elements"][j]
                    #             slides[i]["elements"][j]['x']=data.x#['x']
                    #             slides[i]["elements"][j]['y'] = data.y#['y']
                    #             slides[i]["elements"][j]['width'] = data.cx#['cx']
                    #             slides[i]["elements"][j]['height'] = data.cy#['cy']
                    #
                    #             params = elements["action"]["params"]
                    #             # add new image filename
                    #             #---------------------------
                    #             params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
                    #             #---------------------------
                    #             # add random uuid
                    #             elements["action"]["subContentId"] = str(uuid.uuid4())
                    #             # set width & height
                    #             # params["file"]["width"] = img_width
                    #             # params["file"]["height"] = img_height
                    # print('#-------------------------end text obj.x.y-------------------')
                if(False):
                    for i, image in enumerate(images):
                        slides = content["presentation"]["slides"]
                        # clone first element
                        # if i > 0 and len(slides) <= i:
                        #     slides.append(deepcopy(slides[0]))
                        elements = slides[i]["elements"][0]
                        params = elements["action"]["params"]
                        # add new image filename
                        params["file"]["path"] = "images/" + image
                        # add random uuid
                        elements["action"]["subContentId"] = str(uuid.uuid4())
                        #set width & height
                        params["file"]["width"] = img_width
                        params["file"]["height"] = img_height
                        ratio = img_width / img_height
                        if ratio > target_ratio:  # wider, need to shrink y
                            elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                            elements["height"] = 100# * target_ratio / ratio
                        elif ratio < target_ratio:  # higher, need to shrink x
                            elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                            elements["width"] = 100 #* ratio / target_ratio
            # save file
            #content =


            zout.writestr("content/content.json",  create_newjson_withoutslash(content))
            #---------♥♥♥♥♥♥♥♥---
            #test add json có gì sửa lại ở đây
            #zout.writestr("content/content.json", json.dumps(content))
            #---------




            # add image files to tip
            for image in images:
                zout.write(os.path.join(image_folder, image), "content/images/" + image)

            # change presentation title
            with zin.open("h5p.json", "r") as h5p:
                content = json.load(h5p)
                content["title"] = title
            zout.writestr("h5p.json", json.dumps(content))
def add_to_json_img_text_newver1(newfile, image_folder, images, title, arrayobj,arrobj_shape,images_media,folder_images_pptx_media):
    print('images folder',image_folder)
    print(folder_images_pptx_media)
    print(images_media)
    print(title)
    print(images)
    exclude_files = ['content/content.json', 'h5p.json']
    exclude_filess = ['content/content.json']

    if getattr(sys, 'frozen', False):
        basedir = sys._MEIPASS
    else:
        basedir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(basedir, "template_new.h5p")

    img_width, img_height = get_image_size(os.path.join(image_folder, images[0]))

    with ZipFile(template, "r") as zin:
        with ZipFile(newfile, "w") as zout:
            # copy all other files
            for item in zin.infolist():
                print("item.filename",item.filename)
                if item.filename not in exclude_files:
                    print(zin.read(item.filename))
                    zout.writestr(item, zin.read(item.filename))

            # add image filenames to content.json
            with zin.open("content/content.json") as fp:
                content = json.load(fp)

                #----------------------
                print('value content is:',content)
                print('value content dump:',json.dumps(content, indent=1))
                # ----------------------
                # for ii in content:
                #     print('value ii',ii)
                #print('conten',content)
                slides = content["presentation"]["slides"]
                print(f"adding images from {image_folder}: {images}")
                for i, image in enumerate(images):
                    #slides = content["presentation"]["slides"]
                    # clone first element
                    if i > 0 and len(slides) <= i:
                        slides.append(deepcopy(slides[0]))
                    elements = slides[i]["elements"][0]
                    params = elements["action"]["params"]
                    # add new image filename
                    params["file"]["path"] = "images/" + image
                    # add random uuid
                    elements["action"]["subContentId"] = str(uuid.uuid4())
                    #set width & height
                    params["file"]["width"] = img_width
                    params["file"]["height"] = img_height
                    ratio = img_width / img_height
                    if ratio > target_ratio:  # wider, need to shrink y
                        elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                        elements["height"] = 100# * target_ratio / ratio
                    elif ratio < target_ratio:  # higher, need to shrink x
                        elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                        elements["width"] = 100 #* ratio / target_ratio


                # #god bles you ♥
                a = [i for i, x in enumerate(arrayobj) if not x]
                for x in a:
                    print('values arrayobj none:', x)







                #add shape obj h5p#god bles u ♥!
                for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrobj_shape)):
                    if i in a:
                        print('slide ',i,'is none')
                    else:
                        print('slide ',i)
                        #print('values slide 0 element 1',slides[0]["elements"][1])
                        for j,data in enumerate(arr_obj):
                            print('value j',j)
                            #and j <= count_keys(data)
                            #if(j>=0):
                            slide["elements"].append(data)
                #add elements slide
                for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                    if i in a:
                        print('slide ',i,'is none')
                    else:
                        print('slide ',i)
                        #print('values slide 0 element 1',slides[0]["elements"][1])
                        for j,data in enumerate(arr_obj):
                            print('value j',j)
                            #and j <= count_keys(data)
                            #if(j>=0):
                            slide["elements"].append(data)
                            #print('check type slide element:',type(slide["elements"]))
                delete_text_element1_inslide(content)
















                            #print('check type slide element:',type(slide["elements"]))

                # delete elements neu can
                #delete_text_element1_inslide(content)
                # for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                #     if i in a:
                #         print('slide ',i,'is none')
                #     else:
                #         print('slide ',i)
                #         for j,data in enumerate(arr_obj):
                #             slide["elements"].append(slides[0]["elements"][1])
                #     #del slides[i]["elements"][1]



                print('#------------------------start with add elemtens--------')
                # list01 = [1, 2, 3]
                # list02 = ['x','y']
                # for i, (a, b) in enumerate(zip(list01, list02)):
                #     print (i, a, b)

                #add_rong_elements_0(arrayobj)
                print('#-------------- start check   arr i-----------')
                for i,data in enumerate(arrayobj):
                    print('obj ',i)
                    for j,dataa in enumerate(data):

                        print('#--------------arr j-----------')
                        #,dataa.y,dataa.cx,dataa.cy,dataa.text
                        #print(dataa.x,dataa.y,dataa.cx,dataa.cy,dataa.text)
                        print('elemsnts',j,data)
                        print('#--------------arr j-----------')
                print('#--------------check  sau khi add rong end arr i-----------')
                #edit_elements_slide(content,arrayobj)

                if(False):
                    for i, (slide,arr_obj) in enumerate(zip(content["presentation"]["slides"], arrayobj)):
                        print('#------------------------start with h5p add text')
                        if(True==True):
                            if i not in a:
                                print('slide', i)
                                #print("gia tri cua arrayobj",i,arr_obj)
                                #print('#-------------------------text obj.x.y-------------------')

                                #get element to add param text to json
                                #edit element slides ♥ god bles u
                                #tính lại x,y, cx,cy
                                #viet method
                                #convert_donvi_text_xycxcy_2_h5p()
                                #print('#-------------------------end text obj.x.y-------------------')
                                #for j,data in enumerate(arr_obj):

                                for j, (elements, data) in enumerate(zip(slide['elements'], arr_obj)):
                                    #add element text cho tung slide


                                    # tính lại x,y, cx,cy
                                    #         viet method
                                    if(True):
                                        if data=='':
                                            print('data',j,' is none\n')
                                        else:

                                            #----------------------------------------- phai xu ly duoc list index
                                            print('#-----------start convert x y cx cy---------------')
                                            print('gia tri array[',i,'],[',j,']')

                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                                            print(data.x, data.y,data.cx,data.cy,data.text)
                                            print('slide[',i,']','elements[',j,']',elements)
                                            print('#-----------end convert x y cx cy-----------------')

                                            #slides[0]["elements"][1]
                                            #elements = slides[i]["elements"][j+1]
                                            #checkslide_arr(elements,data)
                                            if(True):
                                                elements['x']=data.x#['x']
                                                elements['y'] = data.y#['y']
                                                elements['width'] = data.cx#['cx']
                                                elements['height'] = data.cy#['cy']

                                                params = elements["action"]["params"]
                                                # add new image filename
                                                #---------------------------
                                                #params["text"] = replacestring(cpmverttext_h5p(cpmverttext_h5p(data.text)))
                                                params.update({"text": cpmverttext_h5p(data.text)})
                                                #---------------------------
                                                # add random uuid
                                                elements["action"]["subContentId"] = str(uuid.uuid4())
                                            #set width & height
                                            # params["file"]["width"] = img_width
                                            # params["file"]["height"] = img_height
                                    print('#-------------------------end text obj.x.y-------------------')
                    print('#------------------------start with h5p add text')

                                #else:
                                #xx=xx-1
                            #xx=xx+1
                        #

                    # for i, arr_obj in enumerate(arrayobj):
                    #     for j, data in enumerate(arr_obj):
                    #
                    #         # tính lại x,y, cx,cy
                    #         # viet method
                    #         if(1==1 ):
                    #             data.x,data.y,data.cx,data.cy=convert_donvi_text_xycxcy_2_h5p(data.x,data.y,data.cx,data.cy)
                    #             elements = slides[i]["elements"][j]
                    #             slides[i]["elements"][j]['x']=data.x#['x']
                    #             slides[i]["elements"][j]['y'] = data.y#['y']
                    #             slides[i]["elements"][j]['width'] = data.cx#['cx']
                    #             slides[i]["elements"][j]['height'] = data.cy#['cy']
                    #
                    #             params = elements["action"]["params"]
                    #             # add new image filename
                    #             #---------------------------
                    #             params["text"] = cpmverttext_h5p(cpmverttext_h5p(data.text))
                    #             #---------------------------
                    #             # add random uuid
                    #             elements["action"]["subContentId"] = str(uuid.uuid4())
                    #             # set width & height
                    #             # params["file"]["width"] = img_width
                    #             # params["file"]["height"] = img_height
                    # print('#-------------------------end text obj.x.y-------------------')
                if(False):
                    for i, image in enumerate(images):
                        slides = content["presentation"]["slides"]
                        # clone first element
                        # if i > 0 and len(slides) <= i:
                        #     slides.append(deepcopy(slides[0]))
                        elements = slides[i]["elements"][0]
                        params = elements["action"]["params"]
                        # add new image filename
                        params["file"]["path"] = "images/" + image
                        # add random uuid
                        elements["action"]["subContentId"] = str(uuid.uuid4())
                        #set width & height
                        params["file"]["width"] = img_width
                        params["file"]["height"] = img_height
                        ratio = img_width / img_height
                        if ratio > target_ratio:  # wider, need to shrink y
                            elements["y"] = 0#100 * (1 - target_ratio / ratio) / 2
                            elements["height"] = 100# * target_ratio / ratio
                        elif ratio < target_ratio:  # higher, need to shrink x
                            elements["x"] = 0#100 * (1 - ratio / target_ratio) / 2
                            elements["width"] = 100 #* ratio / target_ratio
            # save file
            #content =


            zout.writestr("content/content.json",  create_newjson_withoutslash(content))
            #---------♥♥♥♥♥♥♥♥---
            #test add json có gì sửa lại ở đây
            #zout.writestr("content/content.json", json.dumps(content))
            #---------




            # add image files to tip
            for image in images:
                print('check os.path',os.path.join(image_folder, image))
                zout.write(os.path.join(image_folder, image), "content/images/" + image)
                #zout.write(os.path.join(image_folder, image), "content/imagessss/" + image)
            for image in images_media:
                print('check os.path',os.path.join(image_folder, image))
                print(image[-3:])
                if(   image[-3:]!='wdp'and image[-3:]!='emf'):#image[-3:]!='WAV' and image[-3:]!='wav' and image[-3:]!='wav' and
                    print(image[-3:])
                    zout.write(os.path.join(folder_images_pptx_media, image), "content/images/media/" + image)
                #zout.write(os.path.join(image_folder, image), "content/imagessss/" + image)
#images_media,folder_images_pptx_media
            # change presentation title
            with zin.open("h5p.json", "r") as h5p:
                content = json.load(h5p)
                content["title"] = title
            zout.writestr("h5p.json", json.dumps(content))

def congchuoiobj(arrayobj):
    list_text_slide=[]
    text=''
    for i,data in enumerate(arrayobj):
        print('slide ', i)

        for j,value in enumerate(data):

            print(i,'data',value)
            if(text==''):
                text = text + value
            else:
                text=text+','+value
        list_text_slide.append(text)
        text=''
    return list_text_slide  # ,list_none
def moveshapeto_allslide(filename):
    #for j,da in enumerate(xml_slides):
        #print(j)
    #filename = '10slide.pptx'
    prs = Presentation(filename)

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    print(slides)
    for slide in prs.slides:
        for shape in slide.shapes:
            shape.left = -10000
            shape.top = -10000
    # for i in slides:
    #     xml_slides.remove(i)
    prs.save('1slide.pptx')
if __name__ == "__main__":
    # folder=r'C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p_1_2_deep'
    # title='10slide'
    # newname_fodler_rar, folderr, new_title_rar, ten_folder_rar = edit_name_pptx2rar(title, folder)
    # print('newname_fodler_rar:', newname_fodler_rar)
    # print('new_title_rar:', new_title_rar)
    # print('folderr:', folderr)
    # print('ten_folder_rar:', ten_folder_rar)
    # # --------------------------------------
    # # xử lý lại extra file sau
    # # extra_rarfile(new_title_rar, folderr)
    # # --------------------------------------
    # print('extra rar file')
    # # load_xml_from_folder(os.path.splitext(filepath)[0])
    # # #load link xml
    # name_XMLa, linkfile_xml, xxxxx = load_xml_from_folder(
    #     os.path.join(folder, title))
    # print("loadxml from folder", name_XMLa)
    # print("link file xml", linkfile_xml)
    # print("gia tri:", linkfile_xml[1])
    # print('#------------------------------------')
    # print("xxxxx", xxxxx)
    # print('#------------------------------------')
    # listfiletext = []
    #
    # for i in xxxxx:
    #     listfiletext.append(elementree_withedittext.lay_all_objtext_infilexml(os.path.join(folder, title), i))
    # a = [i for i, x in enumerate(listfiletext) if not x]
    # print('obj None',a)
    # # for i,x in enumerate(listfiletext):
    # #
    # #     if not x:
    # #         print('slide',i,'none')
    # #     else:
    # #         print('slidei', i)
    # #         for j,z in enumerate(x):
    # #             print('element j',j,'z.text',z.text)
    # aaa=dithangobject(listfiletext,a)
    # for i, x in enumerate(aaa):
    #     if not x:
    #         print('slide ',i,'None')
    #     else:
    #         print('slide',i)
    #         for j,data in enumerate(x):
    #             print(j,'obj_text_h5p',data)
    ar=[]
    # x=traveobj_slide('10slide',r'C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p_1_2_deep')
    # x=[[''],['']]
    # for i,data in enumerate(x):
    #     print('slide ', i)
    #
    #     for j,value in enumerate(data):
    #
    #         print(i,'data',value)
    try:
        # Manifest
        print("Powerpoint to h5p Converter.")
        print(f"Version: {VERSION}")
        print(f"Martin Lehmann, {YEAR}")
        print("Licence: BSD-2-Clause")
        print("Source code: https://github.com/MM-Lehmann/pptx2h5p")

        print("develop by tran tuan anh")
        print('he is superman')
        if len(sys.argv) != 2:
            print("Usage : python pptx2h5p.py [file]", file=sys.stderr)
            sys.exit(-1)

        # extract metadata
        filepath = os.path.abspath(sys.argv[1])
        print(f"extracting images from {filepath}.")
        folder = os.path.dirname(filepath)
        filename = os.path.basename(filepath)
        title = os.path.splitext(filename)[0]
        if not os.path.exists(filepath):
            print("No such file!", file=sys.stderr)
            sys.exit(-1)

        # extract images
        #sửa nhè nhẹ link ảnh back ground ở đây.
        newlinkfilee=getbackgroundofslide(filepath)
        ppt2png(newlinkfilee)
        #sua images folder
        #--------------da edit o day ban cu
        # image_folder = os.path.join(folder, title)
        # images = natsorted(os.listdir(image_folder))
        #-----------
        if(True):
            #os.rename(newlinkfilee[:-5],os.path.join(folder, title))
            image_folder = os.path.join(folder, title+str('_cp'))
            images = natsorted(os.listdir(image_folder))
            #images media
            image_media_folder = os.path.join(folder, title + str('_copy\ppt\media'))
            images_media = natsorted(os.listdir(image_media_folder))


            # compile .hp5 file
            newfilename = os.path.splitext(filepath)[0] + ".h5p"
            print("building new "+str(newfilename)+" file")
            print(r" \n gia tri cua title la:"+str(title)+r"\n gia tri cua folder"+str(folder))
            #load xml from da giai nen
            newname_fodler_rar, folderr, new_title_rar, ten_folder_rar = edit_name_pptx2rar(title,folder)
            print('tra ve edit')
            # x = traveobj_slide('10slide', r'C:\Users\Admin\anaconda3\envs\h5p_11\pptx2h5p_1_2_deep')
            # x = [[''], ['']]
            # for i, data in enumerate(x):
            #     print('slide ', i)
            #
            #     for j, value in enumerate(data):
            #         print(i, 'data', value)
            print('newname_fodler_rar:',newname_fodler_rar)
            print('new_title_rar:', new_title_rar)
            print('folderr:', folderr)
            print('ten_folder_rar:', ten_folder_rar)
            #--------------------------------------
            #xử lý lại extra file sau
            #extra_rarfile(new_title_rar, folderr)
            #--------------------------------------
            print('extra rar file')
            cxz,list_shape=traveobj_slide(title, folder)
            # print(os.listdir(folder +str( """_copy\ppt""") + str("""\slides""")))
            print('ket thuc trả ve obj_slide')
            for i,data in enumerate(cxz):
                print(i)
                for j,dataa in enumerate(data):
                    print('dept',dataa.x,dataa.y,dataa.cx,dataa.cy,dataa.text,dataa.color,dataa.size)
            a = [i for i, cxz in enumerate(cxz) if not cxz]
            b = [i for i, list_shape in enumerate(list_shape) if not list_shape]
            print(a)
            #get shape,img in slide
            for i,data in enumerate(list_shape):
                for j,dataa in enumerate(data):
                    print(dataa.x,dataa.y,dataa.cx,dataa.cy)
                    print('check list shape',dataa.img,dataa.color,dataa.prst)
            shape_of_slide_object =list_shape
            #lay object json dang string  add to slide
            print('----start get list_shape_obj_inslide')
            list_shape_obj_ofslide=dithang_shape_object_ver001(shape_of_slide_object,arr_none=b)
            for i,data in enumerate(list_shape_obj_ofslide):
                print('slide list_shape_pbi_ofslide ',i,'\n')
                for j,dataa in enumerate(data):
                    print('check list_shape ',j,dataa,'\n')
            print('----end get list_shape_obj_inslide')
            #get text object in slide
            print('----start get texxt')
            listobj_text_slidee = dithangobject_ver001(cxz, arr_none=a)
            print('----end get list_shape_obj_inslide')
            #get img,shape in slide.
            add_to_json_img_text_newver1(
                    newfilename, image_folder, images, title, arrayobj = listobj_text_slidee,arrobj_shape= list_shape_obj_ofslide,images_media=images_media,folder_images_pptx_media=image_media_folder
                )
            # add_to_json_img_text_newver1(
            #     newfilename, image_folder, images, title, arrayobj=listobj_text_slidee
            #     # arrobj_shape=list_shape_obj_ofslide, images_media=images_media,
            #     # folder_images_pptx_media=image_media_folder
            # )
            # print("Converting successfully finished.")
            input(
                "Press Enter to delete temporary image (export) folder and close this window."
            )
            # zip_name = 'path\to\zip_file'
            # directory_name = 'path\to\directory'
            #
            # # Create 'path\to\zip_file.zip'
            # shutil.make_archive(zip_name, 'zip', directory_name)

            for image in images:
                os.remove(os.path.join(image_folder, image))
            os.rmdir(image_folder)

    except Exception as e:
        print(e, file=sys.stderr)
        input("An error has occured. Press Enter to close this window.")

